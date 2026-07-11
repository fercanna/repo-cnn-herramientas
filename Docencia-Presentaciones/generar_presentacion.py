import os
import sys
import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# --- 1. PARSEO DEL GUIÓN ---
def parsear_guion(ruta_txt):
    """Lee el archivo .txt y extrae las diapositivas de forma más robusta."""
    try:
        with open(ruta_txt, 'r', encoding='utf-8') as f:
            texto = f.read()
            
        patron = r"DIAPOSITIVA (\d+):(.*?)(?=DIAPOSITIVA \d+:|\Z)"
        matches = re.findall(patron, texto, re.DOTALL)
        
        diapositivas = []
        for match in matches:
            numero = int(match[0])
            bloque = match[1].strip()
            
            titulo = bloque.split('\\n')[0].strip()
            
            contenido_match = re.search(r"Contenido Sugerido:(.*?)Guion del Orador:", bloque, re.DOTALL)
            guion_match = re.search(r"Guion del Orador:(.*)", bloque, re.DOTALL)
            
            contenido = contenido_match.group(1).strip() if contenido_match else "Contenido no encontrado."
            guion = guion_match.group(1).strip() if guion_match else "Guion no encontrado."
            
            diapositivas.append({
                "numero": numero,
                "titulo": titulo,
                "contenido": contenido,
                "guion": guion
            })
        return diapositivas
    except Exception as e:
        print(f"❌ Error al leer o parsear el guion: {e}")
        return []

# --- 2. CREACIÓN DE LA PRESENTACIÓN PPTX ESTRUCTURADA ---
def crear_presentacion_estructurada(diapositivas, archivo_salida):
    """Crea un archivo .pptx con contenido en el cuerpo y guion en las notas."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    print(f"\\n🏗️ Creando presentación estructurada en '{archivo_salida}'...")

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Presentación sobre Gestión de Riesgos"
    slide.placeholders[1].text = "Generado por Asistente Gemini para CNNIA"

    for d in diapositivas:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = f"Diapositiva {d['numero']}: {d['titulo']}"
        
        body_shape = slide.placeholders[1]
        tf_body = body_shape.text_frame
        tf_body.text = d['contenido'].replace('*','-')

        notes_slide = slide.notes_slide
        tf_notes = notes_slide.notes_text_frame
        tf_notes.text = d['guion']

    try:
        prs.save(archivo_salida)
        print(f"✅ ¡Presentación estructurada guardada exitosamente!")
    except Exception as e:
        print(f"❌ Error al guardar la presentación: {e}")

# --- 3. FLUJO PRINCIPAL ---
def main():
    parser = argparse.ArgumentParser(description="Generador de Presentaciones Estructuradas")
    parser.add_argument("archivo", help="Ruta al archivo .txt del guion")
    args = parser.parse_args()

    output_dir = Path("Salidas_Generadas")
    output_dir.mkdir(exist_ok=True)
    output_filename = output_dir / "Presentacion_Estructurada_Riesgos.pptx"

    print(f"🚀 Iniciando proceso para: {args.archivo}")
    
    diapositivas = parsear_guion(args.archivo)
    if not diapositivas:
        print("❌ No se encontraron diapositivas válidas en el guion.")
        return

    crear_presentacion_estructurada(diapositivas, output_filename)

    print(f"\\n✅ Proceso completado. Tu archivo está en: {output_filename}")

if __name__ == "__main__":
    main()
