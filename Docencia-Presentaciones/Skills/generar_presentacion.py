import os
import sys
import argparse
import re
from pathlib import Path
from dotenv import load_dotenv
import google.genai as genai
from pptx import Presentation
from pptx.util import Inches

# --- CONFIGURACIÓN INICIAL ---
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")

if not api_key:
    print("❌ ERROR: No se encontró GOOGLE_API_KEY en el archivo .env")
    sys.exit(1)

# No se usa genai.configure() en la nueva librería. La clave se usa al crear un cliente si es necesario.

# --- 1. PARSEO DEL GUIÓN ---
def parsear_guion(ruta_txt):
    """Lee el archivo .txt y extrae las diapositivas con el formato DIAPOSITIVA N:"""
    try:
        texto = Path(ruta_txt).read_text(encoding="utf-8")
        bloques = re.split(r"DIAPOSITIVA \d+:", texto)[1:] 
        diapositivas = []
        
        for i, bloque in enumerate(bloques, 1):
            titulo = bloque.strip().split('\n')[0].strip()
            
            contenido_sugerido_match = re.search(r"Contenido Sugerido:(.*?)Guion del Orador:", bloque, re.DOTALL)
            guion_orador_match = re.search(r"Guion del Orador:(.*)", bloque, re.DOTALL)

            contenido_sugerido = contenido_sugerido_match.group(1).strip() if contenido_sugerido_match else ""
            guion_orador = guion_orador_match.group(1).strip() if guion_orador_match else ""
            
            visual_match = re.search(r"\[Visual\]:(.*)", contenido_sugerido, re.DOTALL)
            visual = visual_match.group(1).strip() if visual_match else "No se especificó visualización."

            diapositivas.append({
                "numero": i,
                "titulo": titulo,
                "contenido": contenido_sugerido,
                "visual": visual,
                "guion": guion_orador
            })
        return diapositivas
    except Exception as e:
        print(f"❌ Error al leer el guion: {e}")
        return []

# --- 2. GENERACIÓN DE IMÁGENES (SIMULADO) ---
def generar_imagen_cnnia(diapositiva):
    """Simula la creación de la visual"""
    estilo_base = "Estilo CNNIA: profesional, minimalista, tecnológico, red de nodos, azul oscuro y gris."
    prompt_final = f"{estilo_base} Imagen para: {diapositiva['visual']}"
    
    print(f"🎨 [SIMULADO] Generando imagen para Slide {diapositiva['numero']}: {diapositiva['titulo']}...")
    # Se comenta la llamada a la API ya que es una simulación y la sintaxis cambió.
    # try:
    #     client = genai.Client(api_key=api_key)
    #     model = client.get_model('gemini-1.5-flash')
    #     return True
    # except Exception as e:
    #     print(f"⚠️ Error en la configuración de la API de GenAI: {e}")
    #     return False
    return True


# --- 3. CREACIÓN DE LA PRESENTACIÓN PPTX ---
def crear_presentacion_pptx(diapositivas, archivo_salida):
    """Crea un archivo .pptx a partir de la lista de diapositivas."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    print(f"\n🏗️ Creando presentación en '{archivo_salida}'...")

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Presentación sobre Gestión de Riesgos"
    subtitle.text = f"Generado para CNNIA"

    for d in diapositivas:
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = f"Diapositiva {d['numero']}: {d['titulo']}"
        
        tf = body.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = "CONTENIDO Y VISUAL:"
        p.font.bold = True

        for line in d['contenido'].split('\n'):
            p = tf.add_paragraph()
            p.text = line.strip()
            p.level = 1

        p = tf.add_paragraph()
        p.text = "\n--- GUION DEL ORADOR ---"
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = d['guion']
        p.level = 1

    try:
        prs.save(archivo_salida)
        print(f"✅ ¡Presentación guardada exitosamente!")
    except Exception as e:
        print(f"❌ Error al guardar la presentación: {e}")

# --- 4. FLUJO PRINCIPAL ---
def main():
    parser = argparse.ArgumentParser(description="Skill de Presentaciones Automáticas CNNIA")
    parser.add_argument("archivo", help="Ruta al archivo .txt del guion")
    args = parser.parse_args()

    output_dir = Path("Salidas_Generadas")
    output_dir.mkdir(exist_ok=True)
    output_filename = output_dir / "Presentacion_Riesgos_ANOA.pptx"

    print(f"🚀 Iniciando proceso de automatización para: {args.archivo}")
    
    diapositivas = parsear_guion(args.archivo)
    if not diapositivas:
        print("❌ No se encontraron diapositivas válidas.")
        return

    for d in diapositivas:
        generar_imagen_cnnia(d)

    crear_presentacion_pptx(diapositivas, output_filename)

    print(f"\n✅ Proceso completado. Tu archivo está en: {output_filename}")

if __name__ == "__main__":
    main()
