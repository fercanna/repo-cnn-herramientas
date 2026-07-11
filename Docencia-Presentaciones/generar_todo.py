import os
import sys
import re
import io
from pathlib import Path
from dotenv import load_dotenv
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

# --- CONFIGURACIÓN ---
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")

if not api_key:
    print("❌ ERROR: No se encontró GOOGLE_API_KEY")
    sys.exit(1)

client = genai.Client(api_key=api_key)

# --- 1. PARSEO DEL GUIÓN ---
def parsear_guion(ruta_txt):
    try:
        texto = Path(ruta_txt).read_text(encoding="utf-8")
        bloques = re.split(r"DIAPOSITIVA \d+:", texto)[1:] 
        diapositivas = []
        
        for i, bloque in enumerate(bloques, 1):
            lineas = bloque.strip().split('\n')
            titulo = lineas[0].strip()
            
            contenido_match = re.search(r"Contenido Sugerido:(.*?)Guion del Orador:", bloque, re.DOTALL)
            guion_match = re.search(r"Guion del Orador:(.*)", bloque, re.DOTALL)
            visual_match = re.search(r"\[Visual\]:(.*)", bloque, re.DOTALL)

            visual_raw = visual_match.group(1).strip().split('\n')[0] if visual_match else "Technology network"
            
            diapositivas.append({
                "numero": i,
                "titulo": titulo,
                "contenido": contenido_match.group(1).strip() if contenido_match else "",
                "visual_es": visual_raw,
                "guion": guion_match.group(1).strip() if guion_match else ""
            })
        return diapositivas
    except Exception as e:
        print(f"❌ Error al leer el guion: {e}")
        return []

# --- 2. OPTIMIZACIÓN DE PROMPT (GEMINI 3.1 LITE) ---
def optimizar_prompt(visual_es):
    instruccion = (
        f"Create a professional image generation prompt in English for: '{visual_es}'. "
        f"Style: CNNIA corporate, minimalist, high-tech, network of nodes and thin glowing lines. "
        f"Colors: Deep navy blue and charcoal grey. 8k resolution, clean composition, sharp focus."
    )
    try:
        response = client.models.generate_content(
            model='gemini-3.1-flash-lite-preview',
            contents=instruccion
        )
        return response.text.strip()
    except:
        return f"Minimalist tech network, blue and grey, {visual_es}"

# --- 3. GENERACIÓN DE IMAGEN (IMAGEN 4) ---
def generar_imagen(prompt_en, numero_slide, output_dir):
    print(f"🎨 Generando imagen para Slide {numero_slide}...")
    try:
        response = client.models.generate_images(
            model='imagen-4.0-generate-001',
            prompt=prompt_en,
            config=types.GenerateImagesConfig(number_of_images=1, output_mime_type='image/png')
        )
        img_bytes = response.generated_images[0].image.image_bytes
        img_path = output_dir / f"slide_{numero_slide:02d}.png"
        img_path.write_bytes(img_bytes)
        return img_path
    except Exception as e:
        print(f"⚠️ Error generando imagen {numero_slide}: {e}")
        return None

# --- 4. ENSAMBLAJE PPTX ---
def crear_pptx_final(diapositivas, archivo_salida, img_dir):
    prs = Presentation()
    
    # Slide de Título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Gestión de Riesgos y Mejora Continua"
    slide.placeholders[1].text = "Consultoría CNNIA - Automatización con IA"

    for d in diapositivas:
        # Layout con título y cuerpo
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = d['titulo']
        
        # Texto del contenido
        body = slide.placeholders[1]
        body.text = d['contenido']
        
        # Insertar imagen si existe
        img_path = img_dir / f"slide_{d['numero']:02d}.png"
        if img_path.exists():
            # Posicionar imagen a la derecha (ajuste manual de ejemplo)
            left = Inches(6)
            top = Inches(1.5)
            width = Inches(3.5)
            slide.shapes.add_picture(str(img_path), left, top, width=width)

        # Notas del orador
        slide.notes_slide.notes_text_frame.text = d['guion']

    prs.save(archivo_salida)
    print(f"✅ Presentación final guardada en: {archivo_salida}")

# --- FLUJO PRINCIPAL ---
def main():
    guion_path = "Salidas_Generadas/Guio_Presentación Gestión de Riesgos_ANOA.txt"
    out_dir = Path("Salidas_Generadas")
    img_dir = out_dir / "imagenes_slides"
    img_dir.mkdir(exist_ok=True)
    
    print(f"🚀 Iniciando flujo completo para: {guion_path}")
    
    diapositivas = parsear_guion(guion_path)
    if not diapositivas: return

    for d in diapositivas:
        prompt_en = optimizar_prompt(d['visual_es'])
        generar_imagen(prompt_en, d['numero'], img_dir)

    crear_pptx_final(diapositivas, out_dir / "Presentacion_Final_Automatizada.pptx", img_dir)

if __name__ == "__main__":
    main()
